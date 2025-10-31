from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


# Helper to add slide with title and content lines and speaker notes
def add_slide(title, lines, notes=None):
    """Helper function to add a slide with consistent formatting"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Format title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    paragraph = title_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue

    # Format content
    content = slide.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()

    for i, line in enumerate(lines):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(28)
        p.space_after = Pt(20)

    # Add notes if provided
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    return slide


# Slide contents based on the cleaned deck provided by the user
slides = [
    (
        "Comparative analysis of business functions in Access Bank Plc and General Electric",
        [
            "Marketing Production Finance Human Resources",
            "Service banking and manufacturing comparison",
            "Prepared by Richmond Olorunsogo",
        ],
        "Title slide. Overview of the presentation and author.",
    ),
    (
        "Executive summary",
        [
            "Purpose: Compare how marketing production finance and human resources operate in Access Bank Plc and General Electric",
            "Key finding: Access Bank uses marketing and finance to drive customer expansion digital transformation and regulatory compliance while GE uses production and engineering led operations supported by finance and HR to deliver capital intensive manufacturing performance",
            "Slides present functional roles strategic links and recommendations",
        ],
        "High level summary for busy readers.",
    ),
    (
        "Methodology and sources",
        [
            "Approach: Comparison by function using recent company reports investor presentations and earnings releases",
            "Frameworks: Value chain Porter Balanced scorecard Resource based view Lean manufacturing and service operations theory",
            "Focus on last two years of performance and strategic messaging",
        ],
        "Explain methodology and conceptual frameworks used in the analysis.",
    ),
    (
        "Company snapshot Access Bank Plc",
        [
            "Business model: Universal bank with retail corporate and commercial segments",
            "Services: Deposits loans payments treasury and digital banking",
            "Strategic emphasis: Digital transformation customer acquisition cross border expansion and regulatory capital management",
        ],
        "Brief profile of Access Bank and strategic focus.",
    ),
    (
        "Company snapshot General Electric",
        [
            "Business model: Industrial manufacturer reorganised into focused businesses such as aerospace energy and health care",
            "Focus: Production scale engineering excellence and long term service contracts",
            "Strategic emphasis: Capital allocation to high return businesses and service aftermarket growth",
        ],
        "Brief profile of GE and strategic focus.",
    ),
    (
        "Overview of the four business functions and why they matter",
        [
            "Marketing shapes demand brand trust segmentation and pricing",
            "Production or operations delivers core product or service through process design quality systems and supply chain",
            "Finance allocates capital manages liquidity risk pricing and performance measurement",
            "Human resources attracts develops retains talent shapes culture and links behaviour to strategy",
        ],
        "Define each function and its strategic importance.",
    ),
    (
        "Marketing in Access Bank",
        [
            "Role: Customer acquisition product bundling digital channels SME and retail segmentation",
            "Strategy link: Marketing informs product design pricing and risk appetite and feeds into budgeted revenue targets",
            "Operational implication: Metrics like customer lifetime value conversion rates and digital adoption inform planning and budgets",
        ],
        "Describe marketing function in the bank context.",
    ),
    (
        "Marketing in General Electric",
        [
            "Role: Demand generation for high value engineered solutions account driven marketing supported by product management and engineering",
            "Strategy link: Marketing supports long sales cycle proposals and service contracts",
            "Operational implication: Inputs feed backlog forecasts production scheduling and capital planning",
        ],
        "Describe marketing in an industrial context.",
    ),
    (
        "Production or operations in Access Bank",
        [
            "Role: Transaction processing payments compliance and platform uptime",
            "Strategy link: Operations influence cost to income ratio and customer satisfaction",
            "Operational implication: Service level monitoring incident management and regulatory reporting integrated into monthly control cycles",
        ],
        "Explain what production means for a service firm.",
    ),
    (
        "Production and manufacturing in General Electric",
        [
            "Role: Capital intensive manufacturing of engines turbines and large equipment",
            "Strategy link: Capacity converts backlog to revenue and influences capital expenditure budgets",
            "Operational implication: Production planning ties to procurement supplier management warranty and service networks",
        ],
        "Explain manufacturing operations and strategic implications.",
    ),
    (
        "Finance in Access Bank",
        [
            "Role: Treasury liquidity management loan loss provisioning and regulatory capital compliance",
            "Strategy link: Finance sets capital allocation limits and pricing for lending and deposits",
            "Operational implication: Budgeting cycles include top down capital planning and bottom up branch level targets with monthly reconciliation",
        ],
        "Describe finance function in the bank and its planning role.",
    ),
    (
        "Finance in General Electric",
        [
            "Role: Capital budgeting project finance investment appraisal transfer pricing and cash flow management",
            "Strategy link: Finance evaluates long lived assets return on invested capital and funds R and D and M and A decisions",
            "Operational implication: Multi year capital plans scenario stress tests and portfolio reviews guide board level allocation",
        ],
        "Describe finance in a capital intensive manufacturer.",
    ),
    (
        "Human resources in Access Bank",
        [
            "Role: Talent acquisition for branch staff digital product teams and compliance professionals",
            "Strategy link: HR supports network expansion retention and upskilling for digital transformation",
            "Operational implication: Coordination with finance for personnel expenses and with operations for service continuity",
        ],
        "Describe HR in the bank and strategic links.",
    ),
    (
        "Human resources in General Electric",
        [
            "Role: Specialist engineering talent apprenticeships leadership development and safety culture",
            "Strategy link: HR supports lean manufacturing and long term capacity building",
            "Operational implication: Workforce planning safety metrics and talent pipelines inform labour cost forecasts",
        ],
        "Describe HR in manufacturing and strategic importance.",
    ),
    (
        "How the functions integrate with corporate planning budgeting and control",
        [
            "Access Bank: Integration aligns to quarterly earnings forecasts regulatory capital constraints and compliance timelines",
            "GE: Production led planning produces multi year backlog based forecasts used for capital expenditure and debt management",
            "Both firms use monthly reporting rolling forecasts and control dashboards to align functions",
        ],
        "Explain integration of functions with planning and control.",
    ),
    (
        "Decision making frameworks used by Access Bank",
        [
            "Frameworks: Scenario planning stress testing regulatory capital and risk appetite frameworks",
            "Planning: Top down strategic planning with bottom up budgeting and monthly rolling forecasts",
            "Governance: Formal risk committees and asset liability committees align financial policy to business targets",
        ],
        "Summarise bank decision making frameworks.",
    ),
    (
        "Decision making frameworks used by General Electric",
        [
            "Frameworks: Portfolio management disciplined capital allocation return on invested capital targets and lean continuous improvement",
            "Planning: Multi year rolling plans with sensitivity analysis for market cycles",
            "Governance: Segment level boards and executive committees set strategy and capital budgets",
        ],
        "Summarise manufacturer decision making frameworks.",
    ),
    (
        "Comparative evaluation of decision making and policy setting",
        [
            "Access Bank emphasises regulatory compliance liquidity deposit growth and risk metrics",
            "GE emphasises engineering returns project NPV free cash flow and backlog conversion",
            "Observation: Service firms integrate regulatory and customer metrics while manufacturers coordinate production and finance over long horizons",
        ],
        "Compare decision making emphases and performance criteria.",
    ),
    (
        "Managerial behaviours in Access Bank",
        [
            "Leadership style: Consensus oriented with emphasis on compliance stakeholder management and relationship banking",
            "Impact: Supports stable growth brand trust and regulatory alignment but can slow fast risk taking",
            "Examples: Board risk committees public investor engagement and governance reporting",
        ],
        "Describe managerial behaviour and its strategic influence.",
    ),
    (
        "Managerial behaviours in General Electric",
        [
            "Leadership style: Engineering and performance driven with emphasis on operational discipline and measurable productivity gains",
            "Impact: Drives performance improvements rapid operational changes and demands technical capability",
            "Examples: Focus on execution backlog management and capital return programmes",
        ],
        "Describe managerial behaviour in GE.",
    ),
    (
        "How managerial behaviours influence strategic outcomes",
        [
            "Access Bank: Risk conscious governance supports deposit mobilisation and trust but constrains speculative investments",
            "GE: Performance culture accelerates productivity gains margin expansion and swift capital allocation but creates strong execution pressure",
            "Sector differences explain why HR and production have different strategic weight",
        ],
        "Synthesize behavioural impacts on strategy.",
    ),
    (
        "Integration differences service versus manufacturing",
        [
            "Service firm emphasis: Marketing and finance interact frequently with customers regulators and digital channels",
            "Manufacturing emphasis: Production engineering and supply chain dominate planning cycles and capital planning",
            "Result: Functional priorities shift with industry logic",
        ],
        "Highlight structural differences between sectors.",
    ),
    (
        "Recommendations and implications for managers",
        [
            "For Access Bank: Strengthen cross functional product squads link incentives to digital adoption and use scenario planning for macro shocks",
            "For GE: Continue investment in production productivity skill development and service aftermarket business and strengthen finance engineering collaboration",
            "Both firms: Use balanced scorecards to align financial operational customer and people metrics",
        ],
        "Actionable recommendations for managers in both firms.",
    ),
    (
        "Conclusion",
        [
            "Summary: The four functions play different but complementary roles in Access Bank and General Electric",
            "Access Bank: Marketing and finance drive customer growth and regulatory compliance while operations and HR support reliability",
            "GE: Production and engineering anchor advantage with finance and HR supporting capital heavy delivery",
            "Decision frameworks differ reflecting sector needs and managerial behaviours shape strategic outcomes",
        ],
        "Concluding synthesis and final remarks.",
    ),
    (
        "References and selected sources",
        [
            "Company annual reports investor presentations and public filings",
            "Academic frameworks including Porter Resource based view Balanced scorecard Kaplan and Norton and lean manufacturing literature",
            "Use company disclosures from the last two years for factual claims",
        ],
        "Reference note for the slide deck user. Detailed citations have been omitted as requested by the user.",
    ),
]

for title, lines, notes in slides:
    add_slide(title, lines, notes)

# Save presentation
file_path = "C:\\Users\\HP\\Searches\\Phython test run\\comparative_analysis_access_ge.pptx"
prs.save(file_path)
print(f"Presentation saved to {file_path}")
